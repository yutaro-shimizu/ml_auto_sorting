import os
from PyPDF2 import PdfReader
from docx import Document
import csv
import pandas as pd
from pptx import Presentation
from xml.etree import ElementTree as ET

def read_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            extracted_text = page.extract_text()
            if extracted_text:
                text += extracted_text
        return {"text": text, "file_name": os.path.basename(file_path), "file_path": file_path}
    except Exception as e:
        return {"text": f"Error reading PDF: {e}", "file_name": os.path.basename(file_path), "file_path": file_path, "error": True}

def read_doc(file_path):
    try:
        doc = Document(file_path)
        text = "\n".join([p.text for p in doc.paragraphs])
        return {"text": text, "file_name": os.path.basename(file_path), "file_path": file_path}
    except Exception as e:
        return {"text": f"Error reading Word document: {e}", "file_name": os.path.basename(file_path), "file_path": file_path, "error": True}

def read_py(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        return {"text": text, "file_name": os.path.basename(file_path), "file_path": file_path}
    except Exception as e:
        return {"text": f"Error reading Python script: {e}", "file_name": os.path.basename(file_path), "file_path": file_path, "error": True}

def read_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        return {"text": text, "file_name": os.path.basename(file_path), "file_path": file_path}
    except Exception as e:
        return {"text": f"Error reading text file: {e}", "file_name": os.path.basename(file_path), "file_path": file_path, "error": True}

def read_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text += " ".join(run.text for run in paragraph.runs) + "\n"
        return {"text": text.strip(), "file_name": os.path.basename(file_path), "file_path": file_path}
    except Exception as e:
        return {"text": f"Error reading PowerPoint file: {e}", "file_name": os.path.basename(file_path), "file_path": file_path, "error": True}

def read_excel(file_path):
    try:
        df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets as a dictionary
        text = ""
        for sheet_name, sheet_data in df.items():
            text += f"Sheet: {sheet_name}\n"
            text += sheet_data.to_string(index=False) + "\n"
        return {"text": text.strip(), "file_name": os.path.basename(file_path), "file_path": file_path}
    except Exception as e:
        return {"text": f"Error reading Excel file: {e}", "file_name": os.path.basename(file_path), "file_path": file_path, "error": True}

def read_csv(file_path):
    try:
        with open(file_path, mode='r', encoding='utf-8') as file:
            reader = csv.reader(file)
            text = "\n".join([", ".join(row) for row in reader])
        return {"text": text.strip(), "file_name": os.path.basename(file_path), "file_path": file_path}
    except Exception as e:
        return {"text": f"Error reading CSV file: {e}", "file_name": os.path.basename(file_path), "file_path": file_path, "error": True}

def read_kml(file_path):
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()

        # Extract data from the KML namespace
        namespace = {"kml": "http://www.opengis.net/kml/2.2"}
        text = ""

        for placemark in root.findall(".//kml:Placemark", namespace):
            name = placemark.find("kml:name", namespace)
            description = placemark.find("kml:description", namespace)
            coordinates = placemark.find(".//kml:coordinates", namespace)
            
            if name is not None:
                text += f"Name: {name.text}\n"
            if description is not None:
                text += f"Description: {description.text}\n"
            if coordinates is not None:
                text += f"Coordinates: {coordinates.text.strip()}\n"
            text += "\n"
        
        return {"text": text.strip(), "file_name": os.path.basename(file_path), "file_path": file_path}
    except Exception as e:
        return {"text": f"Error reading KML: {e}", "file_name": os.path.basename(file_path), "file_path": file_path, "error": True}

def read_image(file_path):
    try:
        if not os.path.exists(file_path):
            raise FileNotFoundError("Image file not found.")
        return {"image_path": file_path, "file_name": os.path.basename(file_path)}
    except Exception as e:
        return {"image_path": file_path, "file_name": os.path.basename(file_path), "error": True, "error_msg": str(e)}
